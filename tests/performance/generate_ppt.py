#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate Concise PPT for Group Meeting
DFS FPGA Performance Analysis - Concise Version
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, content_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.space_after = Pt(14)
    
    return slide

def create_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5 * num_rows)).table
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_image_slide(prs, title, image_path, caption=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    if Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), width=Inches(9))
    
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    return slide

def main():
    script_dir = Path(__file__).parent
    plots_dir = script_dir / 'plots'
    output_file = script_dir / 'DFS_FPGA_Performance_Analysis.pptx'
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(prs, 
        "DFS FPGA AES 加速集成",
        "性能分析与预测")
    
    # Slide 2: 工作概述
    create_content_slide(prs, "工作概述", [
        "已完成:",
        "  ✓ FPGA AES-256 加密模块集成到 DFS",
        "  ✓ CPU vs FPGA 性能对比测试",
        "  ✓ SmartNIC 性能预测模型",
        "",
        "进行中:",
        "  ○ FPGA 动态重配置支持多加密算法"
    ])
    
    # Slide 3: 三种方案
    create_table_slide(prs, "三种实现方案",
        ["方案", "实现方式", "特点"],
        [
            ["OpenSSL", "CPU + AES-NI", "通用性强，CPU 占用高"],
            ["FPGA", "Xilinx U200", "加密卸载，PCIe 传输开销"],
            ["SmartNIC", "网卡集成 FPGA", "零拷贝，CPU 完全卸载"]
        ])
    
    # Slide 4: 测试结果 - 综合图
    image_path = str(plots_dir / 'prediction_overview.png')
    create_image_slide(prs, "性能测试结果", image_path,
        "OpenSSL vs FPGA (实测) vs SmartNIC (预测)")
    
    # Slide 5: 加速比
    image_path = str(plots_dir / 'prediction_speedup.png')
    create_image_slide(prs, "加速比分析", image_path,
        "SmartNIC 小文件加速比可达 1.7x")
    
    # Slide 6: 资源使用
    create_two_column_slide(prs, "资源使用对比",
        "CPU 使用率",
        [
            "OpenSSL: 3-4%",
            "FPGA: 2-3%",
            "SmartNIC: <1%"
        ],
        "内存使用率",
        [
            "OpenSSL: 1.5-2%",
            "FPGA: 1.4-1.9%",
            "SmartNIC: <0.3%"
        ])
    
    # Slide 7: 瓶颈分析
    create_content_slide(prs, "瓶颈分析", [
        "延迟分解:",
        "  网络延迟 (85%) + 加密延迟 (10%) + 其他 (5%)",
        "",
        "关键发现:",
        "  • 网络 I/O 是主要瓶颈，非加密计算",
        "  • FPGA AES 吞吐量 8+ GB/s 无法充分发挥",
        "  • SmartNIC 可消除协议栈开销"
    ])
    
    # Slide 8: 性能总结
    create_table_slide(prs, "性能预测总结",
        ["指标", "OpenSSL", "FPGA", "SmartNIC"],
        [
            ["PUT 吞吐量", "基线", "+5-6%", "+15-25%"],
            ["GET 吞吐量", "基线", "-1-2%", "+10-20%"],
            ["CPU 使用率", "3-4%", "2-3%", "<1%"],
            ["延迟降低", "基线", "0-5%", "10-20%"]
        ])
    
    # Slide 9: 未完成工作
    create_content_slide(prs, "进行中: 动态重配置", [
        "目标: FPGA 支持多加密算法动态切换",
        "",
        "计划支持的算法:",
        "  • AES-256 (已完成)",
        "  • SM4 (国密算法)",
        "  • ChaCha20",
        "",
        "技术方案:",
        "  • 部分重配置 (Partial Reconfiguration)",
        "  • 算法切换延迟 < 10ms"
    ])
    
    # Slide 10: 后续计划
    create_content_slide(prs, "后续计划", [
        "1. 完成动态重配置模块",
        "   • SM4/ChaCha20 内核实现",
        "   • 运行时算法切换机制",
        "",
        "2. SmartNIC 原型验证",
        "   • 零拷贝数据路径实现",
        "   • 验证预测模型准确性",
        "",
        "3. 性能优化",
        "   • 批量加密减少 PCIe 开销"
    ])
    
    # Slide 11: Thanks
    create_title_slide(prs, "谢谢！", "欢迎提问")
    
    prs.save(str(output_file))
    print(f"PPT saved to: {output_file}")

if __name__ == '__main__':
    main()
